### Converts multiple choice questions in the Huawei pptx files into text, then json

from pptx import Presentation
import re
from functions import *

questionList = []
optionList = []
answerList = []
currQuestion = ""
currOptions = ""
currAnswer = ""

prs = Presentation("华为传输IE笔试题库.pptx")
count = 0
for slide in prs.slides:
    for shape in slide.shapes:
        count += 1
        if hasattr(shape, "text"):
            slideText = shape.text
            # print(slideText)
            # print("--------------")
            if (re.match(r"[0-9]+\.", slideText)):
                optionIndex = slideText.find("A.")
                if (optionIndex != -1):
                    currQuestion = slideText[0:optionIndex]
                    currOptions = slideText[optionIndex:]
                    questionList.append(currQuestion)
                    optionList.append(currOptions)
            else:
                currAnswer = slideText
                answerList.append(currAnswer)

# for item in answerList:
#     print(item)
# # print("----------")

### Main

QnAList = []
for i in range(0, len(questionList)):
    currQuestion = questionList[i]
    currOptions = optionList[i]
    currAnswer = answerList[i]
    formattedQuestion = "###Question: " + currQuestion + currOptions
    formattedAnswer = "###Answer: " + currAnswer
    textEntry = formattedQuestion + formattedAnswer
    QnAList.append(textEntry)

list2json(QnAList, "华为传输IE笔试题库.json")